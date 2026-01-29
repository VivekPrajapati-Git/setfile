import os
from PyPDF2 import PdfReader
import docxpy
from pptx import Presentation
import pandas as pd

def read_file(path):
    # file_name = path.split('/')
    # file_name = file_name[-1]

    _,extension = os.path.splitext(path)

    if extension == '.pdf':
        reader = PdfReader(path)
        text = reader.pages[0]
        text = text.extract_text()
        return text
    elif extension == '.docx':
        text = docxpy.process(path)
        return text
    elif extension == '.pptx':
        prs = Presentation(path)
        all_text =  []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        all_text.append(paragraph.text)
        text = "\n".join(all_text)
        return text
    elif extension == '.xlsx':
        df = pd.read_excel(path)
        text = ""
        for col in df.columns:
            text += f"{col} \n"
            text += df[col].astype(str).str.cat(sep = "") + "\n"
        return text
    elif extension == '.txt':
        with open(path, "r") as f:
            content = f.read()
            return content
    else:
        return ""